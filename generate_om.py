#!/usr/bin/env python3
"""
SB 1123 Offering Memorandum Generator v4
Reads from XLS financial model — zero independent calculations.

Usage:
  python3 generate_om.py <path_to_xls> [--photos matt.png joe.png]

The XLS is the single source of truth. This script is a presentation layer only.
"""

import openpyxl, sys, os, math
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

# Default asset paths (relative to this script's directory)
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
ASSETS_DIR = os.path.join(SCRIPT_DIR, 'assets')
DEFAULT_MATT_PHOTO = os.path.join(ASSETS_DIR, 'matt_circle.png')
DEFAULT_JOE_PHOTO = os.path.join(ASSETS_DIR, 'joe_circle.png')

# ════════════════════════════════════════════════════════════════
# XLS READER — Maps cell addresses to DEAL config
# ════════════════════════════════════════════════════════════════
def read_xls(path):
    """Read SB 1123 financial model into flat dict."""
    wb = openpyxl.load_workbook(path, data_only=True)
    a = wb['Assumptions']
    o = wb['Outputs']
    cf = wb['Monthly CF']
    btr = wb['BTR Hold']

    def v(sheet, cell, default=0):
        val = sheet[cell].value
        return val if val is not None else default

    # Parse address components from "4430 LINDBLADE AVE, Los Angeles, 90066"
    raw_addr = str(v(a, 'C5', ''))
    parts = [p.strip() for p in raw_addr.split(',')]
    street = parts[0] if parts else raw_addr
    city = parts[1] if len(parts) > 1 else ''
    zipcode = parts[2] if len(parts) > 2 else ''

    d = {}
    # ── Property ──
    d['address'] = street
    d['city'] = city
    d['zip'] = zipcode
    d['state'] = 'CA'
    d['zoning'] = str(v(a, 'C7', 'R1'))
    d['lot_sf'] = float(v(a, 'C8', 0))
    d['lot_width'] = float(v(a, 'C9', 0))
    d['lot_depth'] = float(v(a, 'C10', 0))
    d['slope_pct'] = float(v(a, 'C11', 0))
    d['beds_baths'] = str(v(a, 'C12', ''))
    d['dom'] = int(v(a, 'C13', 0))

    # ── Acquisition ──
    d['asking_price'] = float(v(a, 'C16', 0))

    # ── Development ──
    d['units'] = int(v(a, 'C20', 0))
    d['unit_sf'] = float(v(a, 'C21', 0))
    d['buildable_sf'] = float(v(a, 'C22', 0))
    d['build_cost_psf'] = float(v(a, 'C23', 0))
    d['hard_costs'] = float(v(a, 'C24', 0))
    d['soft_cost_pct'] = float(v(a, 'C25', 0))
    d['soft_costs'] = float(v(a, 'C26', 0))
    d['demo_cost'] = float(v(a, 'C27', 0))
    d['subdivision_cost'] = float(v(a, 'C28', 0))
    d['ae_cost'] = float(v(a, 'C29', 0))
    d['total_dev_costs'] = float(v(a, 'C30', 0))

    # ── Exit ──
    d['exit_psf'] = float(v(a, 'C35', 0))
    d['gross_revenue'] = float(v(a, 'C36', 0))
    d['tx_cost_pct'] = float(v(a, 'C37', 0))
    d['net_sale_proceeds'] = float(v(a, 'C38', 0))

    # ── Timeline ──
    d['predev_months'] = int(v(a, 'G5', 6))
    d['construction_months'] = int(v(a, 'G6', 12))
    d['sale_months'] = int(v(a, 'G7', 6))
    d['hold_months'] = int(v(a, 'G8', 24))

    # ── Capital Structure ──
    d['equity_total'] = float(v(a, 'G14', 0))
    d['debt_total'] = float(v(a, 'G15', 0))
    d['total_project_cost'] = float(v(a, 'G16', 0))
    d['equity_pct'] = float(v(a, 'G17', 0))
    d['interest_rate'] = float(v(a, 'G18', 0))
    d['orig_fee_pct'] = float(v(a, 'G19', 0))
    d['orig_fee_dollars'] = float(v(a, 'G20', 0))
    d['interest_treatment'] = str(v(a, 'G21', 'PIK'))

    # ── Carry ──
    d['prop_tax_rate'] = float(v(a, 'G26', 0))
    d['monthly_tax'] = float(v(a, 'G27', 0))
    d['insurance_annual'] = float(v(a, 'G28', 0))
    d['monthly_insurance'] = float(v(a, 'G29', 0))

    # ── Fees ──
    d['acq_fee_pct'] = float(v(a, 'G33', 0))
    d['acq_fee_dollars'] = float(v(a, 'G34', 0))
    d['asset_mgmt_monthly'] = float(v(a, 'G35', 0))
    d['dev_mgmt_monthly'] = float(v(a, 'G36', 0))
    d['disposition_fee_pct'] = float(v(a, 'G37', 0))
    d['disposition_fee_dollars'] = float(v(a, 'G38', 0))
    d['total_sponsor_fees'] = float(v(a, 'G39', 0))

    # ── Waterfall ──
    d['lp_pref_rate'] = float(v(a, 'C42', 0))  # annual pref rate (e.g. 0.08)
    d['gp_promote_pct'] = float(v(a, 'C43', 0))  # e.g. 0.20
    d['gp_coinvest_pct'] = float(v(a, 'C44', 0))  # e.g. 0.05
    d['lp_promote_pct'] = 1.0 - d['gp_promote_pct']

    # ── BTR ──
    d['btr_rent_monthly'] = float(v(a, 'C48', 0))
    d['btr_vacancy'] = float(v(a, 'C49', 0))
    d['btr_opex_ratio'] = float(v(a, 'C50', 0))
    d['btr_cap_rate'] = float(v(a, 'C51', 0))
    d['btr_refi_ltv'] = float(v(a, 'C52', 0))
    d['btr_perm_rate'] = float(v(a, 'C53', 0))
    d['btr_rent_growth'] = float(v(a, 'C54', 0))

    # ── From Outputs sheet (pre-computed) ──
    d['lp_moic'] = float(v(o, 'C5', 0))
    d['lp_irr'] = float(v(o, 'C6', 0))
    d['lp_total_dist'] = float(v(o, 'C7', 0))
    d['lp_equity_in'] = float(v(o, 'C8', 0))
    d['lp_net_profit'] = float(v(o, 'C9', 0))
    d['project_margin'] = float(v(o, 'C11', 0))
    d['project_moic'] = float(v(o, 'C12', 0))
    d['all_in_psf'] = float(v(o, 'C24', 0))
    d['gp_promote_dollars'] = float(v(o, 'F9', 0))
    d['gp_total_income'] = float(v(o, 'F11', 0))
    d['gp_fee_load'] = float(v(o, 'F13', 0))

    # ── From Monthly CF (waterfall detail) ──
    d['loan_repayment'] = float(v(cf, 'C33', 0))
    d['net_distributable'] = float(v(cf, 'C34', 0))
    d['lp_roc'] = float(v(cf, 'C37', 0))
    d['gp_roc'] = float(v(cf, 'C38', 0))
    d['profit_after_roc'] = float(v(cf, 'C39', 0))
    d['lp_pref_dollars'] = float(v(cf, 'C40', 0))
    d['remaining_after_pref'] = float(v(cf, 'C41', 0))
    d['lp_share_remaining'] = float(v(cf, 'C43', 0))
    d['gp_coinvest_equity'] = float(v(cf, 'C7', 0))
    d['loan_draws'] = float(v(cf, 'C8', 0))
    d['total_interest'] = float(v(cf, 'C29', 0))
    d['total_prop_tax'] = float(v(cf, 'C19', 0))
    d['total_insurance'] = float(v(cf, 'C20', 0))
    d['total_asset_mgmt'] = float(v(cf, 'C21', 0))
    d['total_dev_mgmt'] = float(v(cf, 'C22', 0))

    # ── From BTR sheet ──
    d['btr_gpi'] = float(v(btr, 'C8', 0))
    d['btr_egi'] = float(v(btr, 'C10', 0))
    d['btr_noi'] = float(v(btr, 'C12', 0))
    d['btr_stabilized_value'] = float(v(btr, 'C17', 0))
    d['btr_effective_ltv'] = float(v(btr, 'C20', 0))
    d['btr_refi_loan'] = float(v(btr, 'C21', 0))
    d['btr_annual_ds'] = float(v(btr, 'C23', 0))
    d['btr_dscr'] = float(v(btr, 'C24', 0))
    d['btr_annual_cf'] = float(v(btr, 'C32', 0))
    d['btr_coc'] = float(v(btr, 'C33', 0))
    d['btr_yoc'] = float(v(btr, 'C34', 0))

    # ── Derived (simple, not financial modeling) ──
    d['break_even_psf'] = d['total_project_cost'] / (d['buildable_sf'] * (1 - d['tx_cost_pct'])) if d['buildable_sf'] > 0 else 0
    d['lot_per_unit'] = d['lot_sf'] / d['units'] if d['units'] > 0 else 0
    d['fee_pct_of_cap'] = d['total_sponsor_fees'] / d['total_project_cost'] if d['total_project_cost'] > 0 else 0

    return d


# ════════════════════════════════════════════════════════════════
# DESIGN SYSTEM
# ════════════════════════════════════════════════════════════════
NAVY = RGBColor(0x1E, 0x29, 0x3B)
TEAL = RGBColor(0x0D, 0x94, 0x88)
LTEAL = RGBColor(0x14, 0xB8, 0xA6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
S100 = RGBColor(0xF1, 0xF5, 0xF9)
S200 = RGBColor(0xE2, 0xE8, 0xF0)
S400 = RGBColor(0x94, 0xA3, 0xB8)
S500 = RGBColor(0x64, 0x74, 0x8B)
S600 = RGBColor(0x47, 0x55, 0x69)
S700 = RGBColor(0x33, 0x41, 0x55)
SW = Inches(10); SH = Inches(5.625)
FH = "Georgia"; FB = "Calibri"

def fm(n):
    if abs(n) >= 1e6: return f"${n/1e6:,.1f}M"
    return f"${n:,.0f}"
def fn(n): return f"{n:,.0f}"
def fp(n): return f"{n:.1%}"

def _r(s, x, y, w, h, fill, line=None, lw=0):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line: sh.line.color.rgb = line; sh.line.width = Pt(lw)
    else: sh.line.fill.background()
    return sh

def _t(s, x, y, w, h, text, sz=11, bold=False, color=S600, font=None, align=PP_ALIGN.LEFT):
    tf = s.shapes.add_textbox(x, y, w, h).text_frame
    tf.word_wrap = True; p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = str(text); r.font.size = Pt(sz)
    r.font.bold = bold; r.font.color.rgb = color; r.font.name = font or FB
    return tf

def bg(s, c=WHITE):
    f = s.background.fill; f.solid(); f.fore_color.rgb = c

def hdr(s, title, sub=None):
    _r(s, 0, 0, SW, Inches(0.65), NAVY)
    _r(s, 0, Inches(0.65), SW, Inches(0.04), TEAL)
    tf = s.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(8), Inches(0.5)).text_frame
    tf.word_wrap = True; p = tf.paragraphs[0]
    r = p.add_run(); r.text = title; r.font.size = Pt(17); r.font.color.rgb = WHITE
    r.font.name = FH; r.font.bold = True
    if sub:
        r2 = p.add_run(); r2.text = "   " + sub
        r2.font.size = Pt(10); r2.font.color.rgb = LTEAL; r2.font.name = FB

def ftr(s, d):
    _t(s, Inches(0.5), Inches(5.28), Inches(5), Inches(0.25),
       f"CONFIDENTIAL  |  Lucid Development, LLC  |  February 2026", sz=7, color=S400)

def stat(s, x, y, w, h, val, label, vs=26, vc=NAVY):
    _r(s, x, y, w, h, WHITE, line=S200, lw=0.5)
    _r(s, x, y, Inches(0.06), h, TEAL)
    _t(s, x+Inches(0.15), y+Inches(0.06), w-Inches(0.2), Inches(0.4),
       val, sz=vs, bold=True, color=vc, align=PP_ALIGN.CENTER)
    _t(s, x+Inches(0.15), y+h-Inches(0.28), w-Inches(0.2), Inches(0.22),
       label, sz=7, color=S500, align=PP_ALIGN.CENTER)

def ncard(s, x, y, w, h, num, title, body):
    _r(s, x, y, w, h, WHITE, line=S200, lw=0.5)
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, x+Inches(0.2), y+Inches(0.12), Inches(0.3), Inches(0.3))
    c.fill.solid(); c.fill.fore_color.rgb = TEAL; c.line.fill.background()
    tf = c.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(num); r.font.size = Pt(12); r.font.bold = True
    r.font.color.rgb = WHITE; r.font.name = FB
    _t(s, x+Inches(0.6), y+Inches(0.12), w-Inches(0.75), Inches(0.25), title, sz=10, bold=True, color=NAVY)
    _t(s, x+Inches(0.6), y+Inches(0.4), w-Inches(0.75), h-Inches(0.5), body, sz=8, color=S600)

def tbl(s, x, y, w, rows, cw, rh=Inches(0.32)):
    nr = len(rows); nc = len(rows[0])
    ts = s.shapes.add_table(nr, nc, x, y, w, rh * nr); t = ts.table
    for i, c in enumerate(cw): t.columns[i].width = c
    for ri, rd in enumerate(rows):
        is_h = ri == 0; is_tot = ri == nr-1 and any("Total" in str(c) for c in rd)
        t.rows[ri].height = rh
        for ci, ct in enumerate(rd):
            cell = t.cell(ri, ci); cell.text = ""; cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]; p.space_before = Pt(0); p.space_after = Pt(0)
            r = p.add_run(); r.text = str(ct); r.font.name = FB
            if is_h: r.font.size=Pt(8); r.font.bold=True; r.font.color.rgb=WHITE
            elif is_tot: r.font.size=Pt(8); r.font.bold=True; r.font.color.rgb=NAVY
            else: r.font.size=Pt(8); r.font.color.rgb=S700
            p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.RIGHT
            cell.margin_left=Inches(0.08); cell.margin_right=Inches(0.08)
            cell.margin_top=Inches(0.03); cell.margin_bottom=Inches(0.03)
            tcPr = cell._tc.get_or_add_tcPr()
            sf = tcPr.makeelement(qn('a:solidFill'), {})
            if is_h: v_hex='1E293B'
            elif is_tot: v_hex='E2E8F0'
            elif ri%2==0: v_hex='F8FAFC'
            else: v_hex='FFFFFF'
            sf.append(sf.makeelement(qn('a:srgbClr'),{'val':v_hex})); tcPr.append(sf)
            for bt in ['a:lnL','a:lnR','a:lnT','a:lnB']:
                ln=tcPr.makeelement(qn(bt),{'w':'6350','cap':'flat'})
                sf2=ln.makeelement(qn('a:solidFill'),{})
                sf2.append(sf2.makeelement(qn('a:srgbClr'),{'val':'E2E8F0'}))
                ln.append(sf2); tcPr.append(ln)
    return ts


# ════════════════════════════════════════════════════════════════
# BUILD PRESENTATION
# ════════════════════════════════════════════════════════════════
def build_om(d, matt_photo=None, joe_photo=None):
    """Build full OM from deal data dict."""
    pres = Presentation()
    pres.slide_width = SW; pres.slide_height = SH

    btr_monthly = d['units'] * d['btr_rent_monthly']

    # ── P1: COVER ────────────────────────────────────────────
    s = pres.slides.add_slide(pres.slide_layouts[6]); bg(s, WHITE)
    _r(s, 0, 0, Inches(0.4), SH, NAVY)
    _r(s, Inches(0.4), 0, SW, Inches(0.04), TEAL)
    _t(s, Inches(0.7), Inches(0.3), Inches(4), Inches(0.25),
       "LUCID DEVELOPMENT", sz=10, bold=True, color=TEAL)
    _r(s, Inches(0.7), Inches(1.4), Inches(1.2), Inches(0.04), TEAL)
    _t(s, Inches(0.7), Inches(1.6), Inches(5.0), Inches(1.0),
       d['address'].upper(), sz=30, bold=True, color=NAVY, font=FH)
    _t(s, Inches(0.7), Inches(2.7), Inches(5.0), Inches(0.3),
       f"{d['city']}, {d['state']} {d['zip']}", sz=14, color=S500)
    _t(s, Inches(0.7), Inches(3.2), Inches(5.0), Inches(0.2),
       "CONFIDENTIAL OFFERING MEMORANDUM", sz=9, bold=True, color=S400)
    _t(s, Inches(0.7), Inches(3.45), Inches(5.0), Inches(0.2),
       f"SB 1123 By-Right Development  |  {d['units']}-Unit Townhome Community", sz=9, color=TEAL)
    mx = Inches(6.3)
    for i, (val, label) in enumerate([
        (fm(d['lp_equity_in']), "Target Equity Raise"),
        (f"{d['lp_moic']:.2f}x", "Target MOIC"),
        (fp(d['lp_irr']), "Target IRR"),
        (f"~{d['hold_months']} mo.", "Target Hold"),
    ]):
        my = Inches(1.2 + i * 0.9)
        _r(s, mx, my, Inches(0.05), Inches(0.65), TEAL)
        _t(s, mx+Inches(0.2), my+Inches(0.02), Inches(3), Inches(0.35), val, sz=20, bold=True, color=NAVY)
        _t(s, mx+Inches(0.2), my+Inches(0.4), Inches(3), Inches(0.2), label, sz=8, color=S500)
    _t(s, Inches(0.7), Inches(5.0), Inches(3), Inches(0.2), "February 2026", sz=9, color=S400)
    _t(s, Inches(5.5), Inches(5.0), Inches(4), Inches(0.2),
       "Prepared by Lucid Development, LLC", sz=9, color=S400, align=PP_ALIGN.RIGHT)

    # ── P2: DISCLAIMER ───────────────────────────────────────
    s = pres.slides.add_slide(pres.slide_layouts[6]); bg(s, WHITE)
    _r(s, 0, 0, SW, Inches(0.04), TEAL)
    _t(s, Inches(1), Inches(0.6), Inches(8), Inches(0.35),
       "CONFIDENTIAL \u2014 NOT FOR DISTRIBUTION", sz=14, bold=True, color=NAVY, font=FH)
    _r(s, Inches(1), Inches(1.0), Inches(2), Inches(0.02), TEAL)
    _t(s, Inches(1), Inches(1.2), Inches(8), Inches(3.5),
       'This Confidential Offering Memorandum has been prepared by Lucid Development, LLC '
       'solely for the purpose of providing prospective investors with information regarding '
       f'a potential investment in {d["address"]}, {d["city"]}, {d["state"]} {d["zip"]}.\n\n'
       'This Memorandum does not constitute an offer to sell or a solicitation of an offer to buy. '
       'Any such offer will be made only via definitive PPM and subscription documents.\n\n'
       'All financial projections, return estimates, and timelines herein represent targeted outcomes '
       'only and are not guarantees. Actual results may differ materially. '
       'Past performance is not indicative of future results.\n\n'
       'Recipients agree not to reproduce, distribute, or disclose without written consent.',
       sz=10, color=S600)

    # ── P3: EXEC SUMMARY ─────────────────────────────────────
    s = pres.slides.add_slide(pres.slide_layouts[6]); bg(s, WHITE); hdr(s, "Executive Summary")
    _t(s, Inches(0.5), Inches(0.85), Inches(5.5), Inches(3.6),
       f"Lucid Development, LLC is seeking approximately {fm(d['lp_equity_in'])} in LP equity "
       f"to acquire, subdivide, and develop {d['address']} \u2014 a {fn(d['lot_sf'])} SF parcel "
       f"in {d['city']}, zoned {d['zoning']} and eligible for ministerial approval under "
       f"California SB 1123.\n\n"
       f"The plan is to acquire for {fm(d['asking_price'])}, demolish the existing structure, "
       f"subdivide into {d['units']} fee-simple parcels, and construct {d['units']} modern "
       f"{d['unit_sf']:,.0f} SF townhomes. Completed units are targeted for sale within "
       f"approximately {d['hold_months']} months.\n\n"
       f"SB 1123 enables this without discretionary approvals, public hearings, or CEQA review. "
       f"A build-to-rent fallback targets {fm(btr_monthly)}/month in gross rental income.",
       sz=9, color=S600)
    for i, (val, label) in enumerate([
        (fm(d['lp_equity_in']), "Target Equity Raise"),
        (fm(d['total_project_cost']), "Total Capitalization"),
        (f"{d['lp_moic']:.2f}x  /  {fp(d['lp_irr'])}", "Target MOIC / IRR"),
        (f"{d['units']} Units  |  ~{d['hold_months']} mo.", "Units / Target Hold"),
    ]):
        stat(s, Inches(6.2), Inches(0.85+i*1.0), Inches(3.3), Inches(0.85), val, label, vs=20)
    ftr(s, d)

    # ── P4: HIGHLIGHTS ────────────────────────────────────────
    s = pres.slides.add_slide(pres.slide_layouts[6]); bg(s, S100); hdr(s, "Investment Highlights")
    ratio = d['exit_psf']/d['all_in_psf'] if d['all_in_psf'] else 0
    for i, (t, b) in enumerate([
        ("By-Right Under SB 1123",
         "Ministerial approval \u2014 no public hearings, no CEQA. 60-day statutory timeline."),
        ("Severe Housing Undersupply",
         f"New townhome inventory in {d['city']} is near zero. SB 1123 unlocks infill supply."),
        (f"{ratio:.1f}x Cost-to-Value Spread",
         f"All-in ~${d['all_in_psf']:,.0f}/SF vs. ${d['exit_psf']:,.0f}/SF exit. "
         f"Break-even at ${d['break_even_psf']:,.0f}/SF."),
        ("Built-in Downside Protection",
         f"BTR fallback targets {fm(btr_monthly)}/mo gross rent, {fm(d['btr_noi'])} annual NOI."),
        ("Institutional-Caliber Sponsors",
         "9 years at Moelis / UBS. $350M+ structured capital. Active ops across LA."),
        ("Aligned Incentives",
         f"GP promote subordinated to {d['lp_pref_rate']:.0%} LP pref. "
         f"Fees ({d['fee_pct_of_cap']:.1%} of cap) below market."),
    ]):
        col = i%2; row = i//2
        ncard(s, Inches(0.5+col*4.7), Inches(0.85+row*1.5), Inches(4.4), Inches(1.3), i+1, t, b)
    ftr(s, d)

    # ── P5: WHY LUCID ─────────────────────────────────────────
    s = pres.slides.add_slide(pres.slide_layouts[6]); bg(s, WHITE); hdr(s, "Why Lucid Development")
    _t(s, Inches(0.5), Inches(1.2), Inches(5.5), Inches(3.2),
       "Matt Lucido and Joe Salerno have spent five years building Yardsworth \u2014 a California "
       "FinTech platform at the intersection of housing policy, capital markets, and residential "
       "development.\n\nThrough Yardsworth, the principals structured $350M+ in equity and credit "
       "term sheets, evaluated hundreds of infill sites, and built proprietary tools for small-lot "
       "subdivision economics.\n\nThey know these deals \u2014 zoning, permitting, construction, "
       "exit math \u2014 because they've studied and prepared for them for five years.",
       sz=10, color=S600)
    for i, (v, l) in enumerate([("5 Years","Co-Founder Partnership"),("$350M+","Capital Structured"),
        ("100+","Infill Sites Evaluated"),("13","Assets Under Management")]):
        stat(s, Inches(6.3), Inches(0.85+i*1.0), Inches(3.2), Inches(0.85), v, l, vs=22)
    ftr(s, d)

    # ── P6: SPONSOR OVERVIEW ──────────────────────────────────
    s = pres.slides.add_slide(pres.slide_layouts[6]); bg(s, WHITE); hdr(s, "Sponsor Overview")
    if matt_photo and os.path.exists(matt_photo):
        s.shapes.add_picture(matt_photo, Inches(0.6), Inches(0.85), Inches(0.75), Inches(0.75))
    _t(s, Inches(1.5), Inches(0.85), Inches(3.3), Inches(0.25), "Matt Lucido", sz=13, bold=True, color=NAVY, font=FH)
    _t(s, Inches(1.5), Inches(1.12), Inches(3.3), Inches(0.2), "Co-Founder & CEO, Yardsworth", sz=8, bold=True, color=TEAL)
    _t(s, Inches(0.6), Inches(1.7), Inches(4.2), Inches(2.0),
       "CEO of Yardsworth, structuring $350M+ in equity and credit for residential housing. "
       "Previously Principal at Wavemaker Partners (global VC). Founded, scaled, and exited RentJolt. "
       "Developed 4 ADUs via structured land leases. Former collegiate athlete, UVA.", sz=8, color=S600)
    if joe_photo and os.path.exists(joe_photo):
        s.shapes.add_picture(joe_photo, Inches(5.3), Inches(0.85), Inches(0.75), Inches(0.75))
    _t(s, Inches(6.2), Inches(0.85), Inches(3.3), Inches(0.25), "Joe Salerno", sz=13, bold=True, color=NAVY, font=FH)
    _t(s, Inches(6.2), Inches(1.12), Inches(3.3), Inches(0.2), "Co-Founder, CFO & CIO, Yardsworth", sz=8, bold=True, color=TEAL)
    _t(s, Inches(5.3), Inches(1.7), Inches(4.2), Inches(2.0),
       "CFO/CIO of Yardsworth. Long-term LA owner-operator: infill, ADUs, adaptive reuse, "
       "duplex conversions, value-add multifamily. 9 years at Moelis & Co. and UBS \u2014 "
       "25+ transactions, many >$1B. MBA honors, UCLA Anderson. MS & BS, USC.", sz=8, color=S600)
    _r(s, Inches(0.5), Inches(3.5), Inches(9.0), Inches(0.03), TEAL)
    _t(s, Inches(0.5), Inches(3.65), Inches(9), Inches(0.2), "PLATFORM & TRACK RECORD", sz=9, bold=True, color=NAVY)
    for i, item in enumerate(["9 rental properties under active management",
        "4 ADU developments (structured land lease model)",
        "$350M+ equity and credit term sheets structured",
        "Active SB 1123 pipeline across Greater LA"]):
        col=i%2; row=i//2; x=Inches(0.5+col*4.7); y=Inches(4.0+row*0.35)
        _r(s, x, y+Inches(0.06), Inches(0.08), Inches(0.08), TEAL)
        _t(s, x+Inches(0.2), y, Inches(4.2), Inches(0.3), item, sz=8, color=S600)
    ftr(s, d)

    # ── P7: SB 1123 ──────────────────────────────────────────
    s = pres.slides.add_slide(pres.slide_layouts[6]); bg(s, WHITE); hdr(s, "California SB 1123", "By-Right Subdivision Framework")
    _t(s, Inches(0.5), Inches(0.85), Inches(9), Inches(0.6),
       "SB 1123 (effective July 1, 2025) allows up to 10 fee-simple townhome lots on qualifying parcels "
       "\u2014 approved ministerially, no public hearing, no CEQA, 60-day approval timeline.", sz=9, color=S600)
    for i,(t,b) in enumerate([("Ministerial Approval","No discretionary review, no CEQA. 60-day statutory."),
        ("Fee-Simple Ownership","Buyers own their lot \u2014 not a condo. Higher $/SF, easier to finance."),
        ("No Affordable Req.","No below-market units for \u226410-unit projects."),
        ("R1 Land Arbitrage","Acquire at SFR pricing, develop at townhome density.")]):
        ncard(s, Inches(0.5+(i%2)*4.7), Inches(1.6+(i//2)*1.5), Inches(4.4), Inches(1.3), i+1, t, b)
    _r(s, Inches(0.5), Inches(4.85), Inches(9), Inches(0.4), NAVY)
    tf=s.shapes[-1].text_frame; p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text="Acquire at R1 Pricing  \u2192  Develop at Townhome Density  \u2192  Sell at New-Con $/SF"
    r.font.size=Pt(10); r.font.color.rgb=LTEAL; r.font.name=FB; r.font.bold=True
    ftr(s, d)

    # ── P8: PROPERTY OVERVIEW ─────────────────────────────────
    s = pres.slides.add_slide(pres.slide_layouts[6]); bg(s, WHITE); hdr(s, "Property Overview")
    props = [
        ("Address", f"{d['address']}, {d['city']}, {d['state']} {d['zip']}"),
        ("Zoning", f"{d['zoning']} \u2014 SB 1123 Eligible"),
        ("Lot Size", f"{fn(d['lot_sf'])} SF ({d['lot_sf']/43560:.2f} acres)"),
        ("Dimensions", f"{d['lot_width']:.0f}' x {d['lot_depth']:.0f}'"),
        ("Asking Price", f"{fm(d['asking_price'])}  (${d['asking_price']/d['lot_sf']:,.0f}/SF lot)"),
        ("Units", f"{d['units']} fee-simple townhomes"),
        ("Unit SF", f"{d['unit_sf']:,.0f} SF habitable"),
        ("Strategy", f"Demolish \u2192 Subdivide \u2192 Build {d['units']} townhomes"),
    ]
    _r(s, Inches(0.4), Inches(0.8), Inches(5.5), Inches(3.6), WHITE, S200, 0.5)
    for i,(k,val) in enumerate(props):
        y=Inches(0.9+i*0.4)
        _t(s, Inches(0.6), y, Inches(1.8), Inches(0.3), k, sz=9, bold=True, color=NAVY)
        _t(s, Inches(2.5), y, Inches(3.2), Inches(0.3), val, sz=9, color=S600)
        if i < len(props)-1: _r(s, Inches(0.6), y+Inches(0.34), Inches(5.1), Inches(0.005), S200)
    # Lot math
    _r(s, Inches(6.2), Inches(0.8), Inches(3.3), Inches(3.6), S100, S200, 0.5)
    _r(s, Inches(6.2), Inches(0.8), Inches(0.06), Inches(3.6), TEAL)
    _t(s, Inches(6.5), Inches(0.95), Inches(2.8), Inches(0.25), "LOT MATH", sz=10, bold=True, color=NAVY)
    for i,(val,l) in enumerate([
        (fn(d['lot_sf'])+" SF", "Total Lot Area"),
        (f"\u00f7 {d['units']} lots", "SB 1123 Subdivision"),
        (f"= {d['lot_per_unit']:,.0f} SF/lot", "Per-Lot Area"),
        (f"{fn(d['buildable_sf'])} SF", "Total New Habitable"),
        (f"{d['unit_sf']:,.0f} SF", "Per-Unit Habitable"),
    ]):
        my=Inches(1.4+i*0.5)
        _t(s, Inches(6.5), my, Inches(2.8), Inches(0.3), val, sz=14, bold=True, color=NAVY)
        _t(s, Inches(6.5), my+Inches(0.25), Inches(2.8), Inches(0.2), l, sz=7, color=S500)
    ftr(s, d)

    # ── P9: CONSTRUCTION BUDGET ───────────────────────────────
    s = pres.slides.add_slide(pres.slide_layouts[6]); bg(s, WHITE); hdr(s, "Construction Budget")
    b_rows = [("Cost Category", "Total", "$/SF", "% of Dev."),
        ("Demolition", fm(d['demo_cost']), f"${d['demo_cost']/d['buildable_sf']:,.0f}" if d['buildable_sf'] else "$0", f"{d['demo_cost']/d['total_dev_costs']:.1%}" if d['total_dev_costs'] else "0%"),
        ("Subdivision / Entitlement", fm(d['subdivision_cost']), f"${d['subdivision_cost']/d['buildable_sf']:,.0f}" if d['buildable_sf'] else "$0", f"{d['subdivision_cost']/d['total_dev_costs']:.1%}" if d['total_dev_costs'] else "0%"),
        ("A&E / Design", fm(d['ae_cost']), f"${d['ae_cost']/d['buildable_sf']:,.0f}" if d['buildable_sf'] else "$0", f"{d['ae_cost']/d['total_dev_costs']:.1%}" if d['total_dev_costs'] else "0%"),
        ("Hard Costs (Vertical)", fm(d['hard_costs']), f"${d['build_cost_psf']:,.0f}", f"{d['hard_costs']/d['total_dev_costs']:.1%}" if d['total_dev_costs'] else "0%"),
        ("Soft Costs ({:.0%} of hard)".format(d['soft_cost_pct']), fm(d['soft_costs']), f"${d['soft_costs']/d['buildable_sf']:,.0f}" if d['buildable_sf'] else "$0", f"{d['soft_costs']/d['total_dev_costs']:.1%}" if d['total_dev_costs'] else "0%"),
        ("Total Development", fm(d['total_dev_costs']), f"${d['total_dev_costs']/d['buildable_sf']:,.0f}" if d['buildable_sf'] else "$0", "100%")]
    tbl(s, Inches(0.5), Inches(0.85), Inches(5.8), b_rows, [Inches(2.4), Inches(1.1), Inches(0.9), Inches(0.9)])
    dev_psf = d['total_dev_costs']/d['buildable_sf'] if d['buildable_sf'] else 0
    for i,(val,l) in enumerate([
        (f"${d['build_cost_psf']:,.0f}/SF","Vertical Cost"),
        (f"${dev_psf:,.0f}/SF","Total Dev Cost"),
        (f"${d['exit_psf']:,.0f}/SF","Exit $/SF (GUI)"),
        (f"${d['break_even_psf']:,.0f}/SF","Break-Even"),
    ]):
        vc=TEAL if "Exit" in l else NAVY
        stat(s, Inches(6.5), Inches(0.85+i*0.95), Inches(3.0), Inches(0.8), val, l, vs=20, vc=vc)
    ftr(s, d)

    # ── P10: COMPS ──────────────────────────────────────────────
    comps = d.get('comps', [])
    comp_label = d.get('comp_label', '')
    s = pres.slides.add_slide(pres.slide_layouts[6]); bg(s, WHITE); hdr(s, "Comparable Sales Analysis")
    sub = comp_label if comp_label else f"Exit $/SF: ${d['exit_psf']:,.0f}"
    _t(s, Inches(0.5), Inches(0.75), Inches(9), Inches(0.25), sub, sz=9, color=S600)
    if comps:
        show = comps[:12]
        c_rows = [("Address", "Sale Price", "$/SF", "SqFt", "Bd/Ba", "Zone", "Yr Built", "Sale Date", "Dist")]
        for c in show:
            c_rows.append((
                c.get('address','')[:28],
                f"${c.get('price',0):,.0f}",
                f"${c.get('ppsf',0):,.0f}",
                f"{c.get('sqft',0):,.0f}",
                f"{c.get('beds','')}/{c.get('baths','')}",
                c.get('zone',''),
                str(c.get('year_built','')),
                c.get('date',''),
                f"{c.get('dist',0):.1f}",
            ))
        prices = [c.get('price',0) for c in show if c.get('price')]
        ppsfs = [c.get('ppsf',0) for c in show if c.get('ppsf')]
        sqfts = [c.get('sqft',0) for c in show if c.get('sqft')]
        c_rows.append((
            f"Average ({len(show)} comps)",
            f"${sum(prices)/len(prices):,.0f}" if prices else "\u2014",
            f"${sum(ppsfs)/len(ppsfs):,.0f}" if ppsfs else "\u2014",
            f"{sum(sqfts)/len(sqfts):,.0f}" if sqfts else "\u2014",
            "", "", "", "", "",
        ))
        crh = Inches(0.22)
        tbl(s, Inches(0.3), Inches(1.0), Inches(9.4), c_rows,
            [Inches(2.2), Inches(1.0), Inches(0.65), Inches(0.65), Inches(0.65),
             Inches(0.6), Inches(0.7), Inches(1.0), Inches(0.6)], rh=crh)
        bar_y = Inches(1.0) + crh * len(c_rows) + Inches(0.1)
    else:
        _r(s, Inches(0.5), Inches(1.0), Inches(9), Inches(2.5), S100, S200, 0.5)
        _r(s, Inches(0.5), Inches(1.0), Inches(0.06), Inches(2.5), TEAL)
        _t(s, Inches(0.8), Inches(1.8), Inches(8.4), Inches(0.6),
           "No comparable sales data available.\nComps require desktop browser with full data loaded.",
           sz=10, color=S500, align=PP_ALIGN.CENTER)
        bar_y = Inches(3.7)
    _r(s, Inches(0.3), bar_y, Inches(9.4), Inches(0.45), NAVY)
    tf=s.shapes[-1].text_frame; p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text=f"Target Exit: ${d['exit_psf']:,.0f}/SF   |   All-In: ${d['all_in_psf']:,.0f}/SF   |   Break-Even: ${d['break_even_psf']:,.0f}/SF"
    r.font.size=Pt(9); r.font.color.rgb=WHITE; r.font.name=FB; r.font.bold=True
    ftr(s, d)

    # ── P11: SOURCES & USES ───────────────────────────────────
    s = pres.slides.add_slide(pres.slide_layouts[6]); bg(s, WHITE); hdr(s, "Sources & Uses of Capital")
    src = [("Source","Amount","% of Total"),
        ("LP Equity", fm(d['lp_equity_in']), fp(d['lp_equity_in']/d['total_project_cost']) if d['total_project_cost'] else "0%"),
        ("GP Co-Invest", fm(d['gp_coinvest_equity']), fp(d['gp_coinvest_equity']/d['total_project_cost']) if d['total_project_cost'] else "0%"),
        ("Construction Loan", fm(d['loan_draws']), fp(d['loan_draws']/d['total_project_cost']) if d['total_project_cost'] else "0%"),
        ("Total Sources", fm(d['total_project_cost']), "100%")]
    tbl(s, Inches(0.5), Inches(0.85), Inches(4.2), src, [Inches(1.6), Inches(1.5), Inches(0.9)])
    carry_total = d['total_prop_tax'] + d['total_insurance']
    use = [("Use","Amount","% of Total"),
        ("Land Acquisition", fm(d['asking_price']), fp(d['asking_price']/d['total_project_cost']) if d['total_project_cost'] else "0%"),
        ("Development Costs", fm(d['total_dev_costs']), fp(d['total_dev_costs']/d['total_project_cost']) if d['total_project_cost'] else "0%"),
        ("Sponsor Fees", fm(d['total_sponsor_fees']), fp(d['total_sponsor_fees']/d['total_project_cost']) if d['total_project_cost'] else "0%"),
        ("Carry (Tax+Ins)", fm(carry_total), fp(carry_total/d['total_project_cost']) if d['total_project_cost'] else "0%"),
        ("Interest (PIK)", fm(d['total_interest']), fp(d['total_interest']/d['total_project_cost']) if d['total_project_cost'] else "0%"),
        ("Origination Fee", fm(d['orig_fee_dollars']), fp(d['orig_fee_dollars']/d['total_project_cost']) if d['total_project_cost'] else "0%"),
        ("Total Uses", fm(d['total_project_cost']), "100%")]
    tbl(s, Inches(5.2), Inches(0.85), Inches(4.3), use, [Inches(1.6), Inches(1.5), Inches(0.9)])
    _r(s, Inches(0.5), Inches(4.2), Inches(9), Inches(0.55), NAVY)
    _r(s, Inches(0.5), Inches(4.2), Inches(0.06), Inches(0.55), TEAL)
    _t(s, Inches(0.8), Inches(4.25), Inches(8.5), Inches(0.45),
       f"All-In: ${d['all_in_psf']:,.0f}/SF  |  Equity: {fm(d['equity_total'])} ({fp(d['equity_pct'])})  |  "
       f"Debt: {fm(d['debt_total'])} at {fp(d['interest_rate'])}  |  Interest: {d['interest_treatment']}",
       sz=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    ftr(s, d)

    # ── P12: PRO FORMA ────────────────────────────────────────
    s = pres.slides.add_slide(pres.slide_layouts[6]); bg(s, WHITE); hdr(s, "Targeted Pro Forma Returns")
    pf = [("Line Item","Amount"),
        (f"Gross Revenue ({d['units']} x {d['unit_sf']:,.0f} SF x ${d['exit_psf']:,.0f}/SF)", fm(d['gross_revenue'])),
        (f"Less: Transaction Costs ({fp(d['tx_cost_pct'])})", f"({fm(d['gross_revenue']-d['net_sale_proceeds'])})"),
        ("Net Sale Proceeds", fm(d['net_sale_proceeds'])),
        ("Less: Loan Repayment (principal + PIK interest)", f"({fm(d['loan_repayment'])})"),
        ("Net Distributable Cash", fm(d['net_distributable'])),
        ("LP Return of Capital", fm(d['lp_roc'])),
        (f"LP Preferred Return ({fp(d['lp_pref_rate'])} annual)", fm(d['lp_pref_dollars'])),
        ("Remaining After Pref", fm(d['remaining_after_pref'])),
        (f"GP Promote ({fp(d['gp_promote_pct'])})", fm(d['gp_promote_dollars'])),
        (f"LP Share of Remaining ({fp(d['lp_promote_pct'])})", fm(d['lp_share_remaining'])),
        ("Total LP Distribution", fm(d['lp_total_dist']))]
    tbl(s, Inches(0.5), Inches(0.85), Inches(5.5), pf, [Inches(4.0), Inches(1.5)])
    stat(s, Inches(6.3), Inches(0.85), Inches(3.2), Inches(0.9), f"{d['lp_moic']:.2f}x", "Target LP MOIC", vs=28, vc=NAVY)
    stat(s, Inches(6.3), Inches(1.95), Inches(3.2), Inches(0.9), fp(d['lp_irr']), "Target LP IRR (XIRR)", vs=28, vc=NAVY)
    stat(s, Inches(6.3), Inches(3.05), Inches(3.2), Inches(0.9), f"~{d['hold_months']} mo.", "Target Hold", vs=28, vc=TEAL)
    stat(s, Inches(6.3), Inches(4.15), Inches(3.2), Inches(0.9), f"{fp(d['lp_pref_rate'])} annual", "LP Preferred Return", vs=28, vc=TEAL)
    _t(s, Inches(0.5), Inches(5.0), Inches(5), Inches(0.25),
       "Projected targets only. Actual results may differ.", sz=7, color=S400)
    ftr(s, d)

    # ── P13: SENSITIVITY (placeholder for XLS) ────────────────
    s = pres.slides.add_slide(pres.slide_layouts[6]); bg(s, WHITE); hdr(s, "Sensitivity Analysis")
    _t(s, Inches(0.5), Inches(0.85), Inches(9), Inches(0.4),
       "Two-way sensitivity tables from the financial model. Rows/columns vary exit $/SF, "
       "build cost $/SF, and hold period to show LP MOIC under different scenarios.", sz=9, color=S600)
    # Table 1 header
    _t(s, Inches(0.5), Inches(1.4), Inches(9), Inches(0.25),
       f"TABLE 1: LP MOIC \u2014 Exit $/SF vs Build Cost $/SF  (Base: ${d['exit_psf']:,.0f}/SF, ${d['build_cost_psf']:,.0f}/SF)",
       sz=9, bold=True, color=NAVY)
    _r(s, Inches(0.5), Inches(1.7), Inches(9), Inches(1.2), S100, S200, 0.5)
    _t(s, Inches(1.5), Inches(2.1), Inches(7), Inches(0.5),
       "Sensitivity values computed in XLS model.\nWill auto-populate when formula engine is connected.",
       sz=9, color=S400, align=PP_ALIGN.CENTER)
    _t(s, Inches(0.5), Inches(3.2), Inches(9), Inches(0.25),
       f"TABLE 2: LP MOIC \u2014 Exit $/SF vs Hold Period  (Base: ${d['exit_psf']:,.0f}/SF, {d['hold_months']} months)",
       sz=9, bold=True, color=NAVY)
    _r(s, Inches(0.5), Inches(3.5), Inches(9), Inches(1.2), S100, S200, 0.5)
    _t(s, Inches(1.5), Inches(3.9), Inches(7), Inches(0.5),
       "Sensitivity values computed in XLS model.\nWill auto-populate when formula engine is connected.",
       sz=9, color=S400, align=PP_ALIGN.CENTER)
    ftr(s, d)

    # ── P14: BTR FALLBACK ─────────────────────────────────────
    s = pres.slides.add_slide(pres.slide_layouts[6]); bg(s, WHITE); hdr(s, "Build-to-Rent Fallback", "Downside Protection")
    _t(s, Inches(0.5), Inches(0.85), Inches(5.3), Inches(2.2),
       "If sale conditions are unfavorable, units convert to a rental portfolio. "
       "Refinance into I/O permanent debt, hold for cash flow + appreciation.\n\n"
       f"With {fp(d['btr_rent_growth'])} annual rent growth, cash-on-cash improves each year. "
       f"DSCR floor of {d.get('btr_dscr',1.25):.2f}x constrains LTV to maintain lender requirements.\n\n"
       "Rent assumptions from GUI deal finder.", sz=9, color=S600)
    for i,(val,l) in enumerate([
        (f"${d['btr_rent_monthly']:,.0f}/mo", "Rent / Unit"),
        (fm(d['btr_gpi']), "Gross Potential Income"),
        (fm(d['btr_noi']), f"NOI ({1-d['btr_opex_ratio']:.0%} margin)"),
        (fm(d['btr_stabilized_value']), f"Stabilized Value ({fp(d['btr_cap_rate'])} cap)"),
        (f"{d['btr_dscr']:.2f}x", "DSCR (I/O)"),
        (fp(d['btr_yoc']), "Yield on Cost"),
    ]):
        stat(s, Inches(6.2), Inches(0.85+i*0.72), Inches(3.3), Inches(0.62), val, l, vs=16)
    _r(s, Inches(0.5), Inches(4.8), Inches(9), Inches(0.45), NAVY)
    _r(s, Inches(0.5), Inches(4.8), Inches(0.06), Inches(0.45), TEAL)
    _t(s, Inches(0.8), Inches(4.85), Inches(8.5), Inches(0.35),
       f"If sale market unfavorable, target rent: {fm(btr_monthly)}/month gross.",
       sz=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    ftr(s, d)

    # ── P15: WATERFALL ────────────────────────────────────────
    s = pres.slides.add_slide(pres.slide_layouts[6]); bg(s, WHITE); hdr(s, "Deal Structure & Waterfall")
    lp_pct_str = fp(d['lp_promote_pct']); gp_pct_str = fp(d['gp_promote_pct'])
    wf = [("Priority","Distribution","Recipient"),
        ("1st","Debt Repayment","Lender"),
        ("2nd","Return of Capital","100% LP + GP (pro rata)"),
        ("3rd",f"LP Preferred ({fp(d['lp_pref_rate'])} annual)","100% LP"),
        ("4th","Remaining Profits",f"{lp_pct_str} LP / {gp_pct_str} GP")]
    tbl(s, Inches(0.5), Inches(0.85), Inches(4.3), wf, [Inches(0.8), Inches(2.0), Inches(1.5)], rh=Inches(0.35))
    fee = [("Fee","Basis","Amount"),
        ("Acquisition Fee", fp(d['acq_fee_pct']), fm(d['acq_fee_dollars'])),
        ("Asset Mgmt", f"${d['asset_mgmt_monthly']:,.0f}/mo", fm(d['total_asset_mgmt'])),
        ("Dev Mgmt", f"${d['dev_mgmt_monthly']:,.0f}/mo (constr)", fm(d['total_dev_mgmt'])),
        ("Disposition Fee", fp(d['disposition_fee_pct']), fm(d['disposition_fee_dollars'])),
        ("Total Fees", fp(d['fee_pct_of_cap']), fm(d['total_sponsor_fees']))]
    tbl(s, Inches(5.2), Inches(0.85), Inches(4.3), fee, [Inches(1.5), Inches(1.2), Inches(1.3)], rh=Inches(0.35))
    _r(s, Inches(0.5), Inches(3.3), Inches(9), Inches(0.5), S100); _r(s, Inches(0.5), Inches(3.3), Inches(0.06), Inches(0.5), TEAL)
    _t(s, Inches(0.8), Inches(3.35), Inches(8.5), Inches(0.4),
       f"GP co-invests {fp(d['gp_coinvest_pct'])} alongside LPs. "
       f"Promote subordinated to {fp(d['lp_pref_rate'])} annual pref. "
       f"Fees ({fp(d['fee_pct_of_cap'])} of cap) below market for ground-up development.",
       sz=9, bold=True, color=NAVY)
    _t(s, Inches(0.5), Inches(4.0), Inches(9), Inches(0.2), "SPONSOR ECONOMICS", sz=9, bold=True, color=NAVY)
    _t(s, Inches(0.5), Inches(4.25), Inches(9), Inches(0.4),
       f"GP Total Income: {fm(d['gp_total_income'])}  |  Promote: {fm(d['gp_promote_dollars'])}  |  "
       f"Fees: {fm(d['total_sponsor_fees'])}  |  GP Fee Load: {fp(d['gp_fee_load'])} of gross revenue",
       sz=8, color=S600)
    ftr(s, d)

    # ── P16: TERMS ────────────────────────────────────────────
    s = pres.slides.add_slide(pres.slide_layouts[6]); bg(s, WHITE); hdr(s, "Investment Terms")
    terms = [("Term","Detail"),
        ("LP Equity", f"~{fm(d['lp_equity_in'])}"),
        ("GP Co-Invest", f"{fp(d['gp_coinvest_pct'])} ({fm(d['gp_coinvest_equity'])})"),
        ("Total Equity", fm(d['equity_total'])),
        ("LP Preferred Return", f"{fp(d['lp_pref_rate'])} annual (cumulative)"),
        ("Profit Split (above pref)", f"{lp_pct_str} LP / {gp_pct_str} GP"),
        ("Target Hold", f"~{d['hold_months']} months"),
        ("Capital Call", "100% at close"),
        ("Distributions", "Upon sale or refinancing event")]
    tbl(s, Inches(0.5), Inches(0.85), Inches(5.5), terms, [Inches(2.2), Inches(3.3)], rh=Inches(0.38))
    _r(s, Inches(6.3), Inches(0.85), Inches(3.2), Inches(4.0), S100, S200, 0.5)
    _r(s, Inches(6.3), Inches(0.85), Inches(0.06), Inches(4.0), TEAL)
    _t(s, Inches(6.6), Inches(1.0), Inches(2.7), Inches(0.25), "GP ALIGNMENT", sz=10, bold=True, color=NAVY)
    for i, item in enumerate([
        f"{fp(d['lp_pref_rate'])} annual pref before GP promote",
        f"GP co-invests {fp(d['gp_coinvest_pct'])} alongside LPs",
        f"Fees at {fp(d['fee_pct_of_cap'])} of cap (below market)",
        "Both principals active: acquisition \u2192 disposition",
        "No GP promote until LP pref is met",
    ]):
        ay = Inches(1.45+i*0.5)
        _r(s, Inches(6.6), ay+Inches(0.04), Inches(0.08), Inches(0.08), TEAL)
        _t(s, Inches(6.85), ay, Inches(2.5), Inches(0.4), item, sz=8, color=S600)
    ftr(s, d)

    # ── P17: RISKS ────────────────────────────────────────────
    s = pres.slides.add_slide(pres.slide_layouts[6]); bg(s, S100); hdr(s, "Risk Factors & Mitigants")
    ratio_str = f"{d['exit_psf']/d['all_in_psf']:.1f}x" if d['all_in_psf'] else "N/A"
    for i,(t,b) in enumerate([
        ("Construction Overruns", f"Soft costs at {fp(d['soft_cost_pct'])} of hard. ${d['all_in_psf']:,.0f}/SF all-in provides {ratio_str} margin."),
        ("Entitlement Risk", "SB 1123 ministerial \u2014 no CEQA, 60-day statutory."),
        ("Market / Pricing", f"Break-even at ${d['break_even_psf']:,.0f}/SF. BTR fallback at {fm(btr_monthly)}/mo."),
        ("Interest Rate / Carry", f"{d['interest_treatment']} interest at {fp(d['interest_rate'])}. Equity sized to maintain LLC cash \u2265 $0."),
        ("Construction Timeline", "Type V wood-frame. Both principals have direct oversight."),
        ("Liquidity", f"Illiquid, ~{d['hold_months']}-month commitment. BTR provides income fallback."),
    ]):
        ncard(s, Inches(0.5+(i%2)*4.7), Inches(0.85+(i//2)*1.5), Inches(4.4), Inches(1.3), i+1, t, b)
    ftr(s, d)

    # ── P18: TIMELINE ─────────────────────────────────────────
    s = pres.slides.add_slide(pres.slide_layouts[6]); bg(s, WHITE); hdr(s, "Project Timeline")
    pm = d['predev_months']; cm_val = d['construction_months']; sm = d['sale_months']; hm = d['hold_months']
    cs = pm  # construction start
    gantt = [
        ("Acquisition & Close", 0, 2, TEAL),
        ("Pre-Dev (Demo/Subdiv/A&E)", 1, pm, RGBColor(0x0D,0x94,0x88)),
        ("Subdivision Map & Approval", 2, pm+cm_val, LTEAL),
        ("Permitting", 2, pm, RGBColor(0x5E,0xEA,0xD4)),
        ("Site Work", cs-1, cs+2, S500),
        ("Vertical Construction", cs, cs+cm_val, NAVY),
        ("Sales & Marketing", cs+cm_val-2, cs+cm_val+sm-2, RGBColor(0x0D,0x94,0x88)),
        ("Final Sales", cs+cm_val+sm-4, hm, TEAL),
    ]
    cx=Inches(2.8); cw=Inches(6.7); cy=Inches(0.95); rh_val=Inches(0.5); mx_val=hm+2
    for m in range(0, mx_val+1, 3):
        px=cx+(m/mx_val)*cw
        _r(s, px, cy, Inches(0.005), rh_val*len(gantt), S200)
        _t(s, px-Inches(0.15), cy-Inches(0.22), Inches(0.35), Inches(0.2), f"M{m}", sz=7, color=S400, align=PP_ALIGN.CENTER)
    for i,(name,start,end,color) in enumerate(gantt):
        y=cy+rh_val*i
        if i%2==0: _r(s, Inches(0.3), y, Inches(9.4), rh_val, S100)
        _t(s, Inches(0.35), y+Inches(0.1), Inches(2.3), Inches(0.3), name, sz=8, color=S700)
        bx=cx+(max(0,start)/mx_val)*cw; bw=max(Inches(0.1),((min(end,mx_val)-max(0,start))/mx_val)*cw)
        bar=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, y+Inches(0.1), bw, Inches(0.25))
        bar.fill.solid(); bar.fill.fore_color.rgb=color; bar.line.fill.background()
    ftr(s, d)

    # ── P19: CTA ──────────────────────────────────────────────
    s = pres.slides.add_slide(pres.slide_layouts[6]); bg(s, WHITE)
    _r(s, 0, 0, Inches(0.4), SH, NAVY); _r(s, Inches(0.4), 0, SW, Inches(0.04), TEAL)
    _r(s, Inches(3.5), Inches(0.8), Inches(3), Inches(0.04), TEAL)
    _t(s, Inches(1), Inches(1.0), Inches(8), Inches(0.5), "Investment Opportunity", sz=28, bold=True, color=NAVY, font=FH, align=PP_ALIGN.CENTER)
    _t(s, Inches(1), Inches(1.55), Inches(8), Inches(0.3), f"{d['address']}  |  {d['city']}, {d['state']}", sz=12, color=S500, align=PP_ALIGN.CENTER)
    for i,(val,l) in enumerate([
        (fm(d['lp_equity_in']),"Target Equity"),(f"{d['lp_moic']:.2f}x","Target MOIC"),
        (fp(d['lp_irr']),"Target IRR"),(f"~{d['hold_months']} mo.","Target Hold"),
    ]):
        sx=Inches(0.8+i*2.3)
        _r(s, sx+Inches(0.35), Inches(2.2), Inches(0.05), Inches(0.7), TEAL)
        _t(s, sx+Inches(0.55), Inches(2.2), Inches(1.6), Inches(0.4), val, sz=22, bold=True, color=NAVY)
        _t(s, sx+Inches(0.55), Inches(2.6), Inches(1.6), Inches(0.25), l, sz=8, color=S500)
    _r(s, Inches(3), Inches(3.8), Inches(4), Inches(0.02), S200)
    _t(s, Inches(1), Inches(4.0), Inches(8), Inches(0.25), "Matt Lucido  |  matt@yardsworth.com", sz=10, color=NAVY, align=PP_ALIGN.CENTER)
    _t(s, Inches(1), Inches(4.25), Inches(8), Inches(0.25), "Joe Salerno  |  joe@yardsworth.com", sz=10, color=NAVY, align=PP_ALIGN.CENTER)
    _t(s, Inches(1), Inches(4.8), Inches(8), Inches(0.25), "Lucid Development, LLC", sz=11, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

    return pres


# ════════════════════════════════════════════════════════════════
# MAIN
# ════════════════════════════════════════════════════════════════
if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python3 generate_om.py <path_to_xls> [--photos matt.png joe.png]")
        sys.exit(1)

    xls_path = sys.argv[1]

    # Parse optional --photos flag
    matt_photo = DEFAULT_MATT_PHOTO
    joe_photo = DEFAULT_JOE_PHOTO
    if '--photos' in sys.argv:
        idx = sys.argv.index('--photos')
        if idx + 2 < len(sys.argv):
            matt_photo = sys.argv[idx + 1]
            joe_photo = sys.argv[idx + 2]

    if not os.path.exists(xls_path):
        print(f"Error: XLS file not found: {xls_path}")
        sys.exit(1)

    print(f"Reading: {xls_path}")
    d = read_xls(xls_path)

    print(f"\n{'='*50}")
    print(f"DEAL: {d['address']}, {d['city']}, {d['state']} {d['zip']}")
    print(f"{'='*50}")
    print(f"Asking:      {fm(d['asking_price'])}")
    print(f"Units:       {d['units']}")
    print(f"Exit $/SF:   ${d['exit_psf']:,.0f}")
    print(f"All-In $/SF: ${d['all_in_psf']:,.0f}")
    print(f"Break-Even:  ${d['break_even_psf']:,.0f}")
    print(f"Total Cap:   {fm(d['total_project_cost'])}")
    print(f"LP Equity:   {fm(d['lp_equity_in'])}")
    print(f"LP MOIC:     {d['lp_moic']:.2f}x")
    print(f"LP IRR:      {d['lp_irr']:.1%}")
    print(f"Waterfall:   {d['lp_pref_rate']:.0%} pref / {d['gp_promote_pct']:.0%} GP promote / {d['gp_coinvest_pct']:.0%} GP co-invest")

    pres = build_om(d, matt_photo, joe_photo)

    # Output PPTX next to the input XLS: {ADDRESS}-OM.pptx
    xls_dir = os.path.dirname(os.path.abspath(xls_path))
    addr_slug = d['address'].replace(' ', '-').replace(',', '').replace('.', '')
    out = os.path.join(xls_dir, f"{addr_slug}-OM.pptx")
    pres.save(out)
    print(f"\nSaved: {out}  ({len(pres.slides)} slides)")
